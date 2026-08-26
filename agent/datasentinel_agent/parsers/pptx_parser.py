"""PPTX parser (python-pptx) — slide-addressable (reported via `page_number`)."""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from datasentinel_agent.config.scan_config import load_scan_config
from datasentinel_agent.parsers.archive_guard import check_zip_safety
from datasentinel_agent.parsers.base import DocumentParser, ExtractedUnit, ParserError


class PPTXParser(DocumentParser):
    def can_parse(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def extract_units(self, file_path: Path) -> Iterator[ExtractedUnit]:
        limits = load_scan_config().archive_limits
        check_zip_safety(
            file_path,
            max_members=limits.max_members,
            max_uncompressed_bytes=limits.max_uncompressed_bytes,
            max_ratio=limits.max_ratio,
        )

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserError("python-pptx is not installed") from exc

        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:
            raise ParserError(f"Failed to open PPTX {file_path}: {exc}") from exc

        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text)
            if texts:
                yield ExtractedUnit(text="\n".join(texts), page_number=slide_number)
