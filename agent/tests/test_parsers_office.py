"""Phase 4 tests for the binary office formats. Fixtures are generated
on-the-fly with the same libraries used to read them (round-trip), keeping
the test suite free of committed binary files and any real personal data."""

import pytest

from datasentinel_agent.parsers.registry import safe_extract


def test_pdf_parser_extracts_page_text(tmp_path):
    fitz = pytest.importorskip("fitz")
    f = tmp_path / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Synthetic record: alex.synthetic@example.com")
    doc.save(str(f))
    doc.close()

    units, error = safe_extract(f)
    assert error is None
    assert units[0].page_number == 1
    assert "alex.synthetic@example.com" in units[0].text


def test_docx_parser_extracts_paragraphs_and_tables(tmp_path):
    docx = pytest.importorskip("docx")
    f = tmp_path / "sample.docx"
    document = docx.Document()
    document.add_paragraph("Employee: Synthetic Person")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "SSN"
    table.rows[0].cells[1].text = "123-45-6789"
    document.save(str(f))

    units, error = safe_extract(f)
    assert error is None
    joined = " ".join(u.text for u in units)
    assert "Synthetic Person" in joined
    assert "123-45-6789" in joined


def test_xlsx_parser_extracts_rows_with_sheet_name(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    f = tmp_path / "sample.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Employees"
    sheet.append(["name", "email"])
    sheet.append(["Synthetic Person", "synthetic.person@example.com"])
    workbook.save(str(f))

    units, error = safe_extract(f)
    assert error is None
    assert all(u.sheet_name == "Employees" for u in units)
    assert any("synthetic.person@example.com" in u.text for u in units)


def test_pptx_parser_extracts_slide_text(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    f = tmp_path / "sample.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = "Contact: synthetic.contact@example.com"
    presentation.save(str(f))

    units, error = safe_extract(f)
    assert error is None
    assert units[0].page_number == 1
    assert "synthetic.contact@example.com" in units[0].text


def test_corrupted_pdf_never_raises(tmp_path):
    pytest.importorskip("fitz")
    f = tmp_path / "corrupt.pdf"
    f.write_bytes(b"%PDF-1.4 this is not a real pdf structure")
    units, error = safe_extract(f)
    assert units == []
    assert error is not None


def test_corrupted_docx_never_raises(tmp_path):
    pytest.importorskip("docx")
    f = tmp_path / "corrupt.docx"
    f.write_bytes(b"not a real zip/docx")
    units, error = safe_extract(f)
    assert units == []
    assert error is not None


def test_corrupted_csv_never_raises(tmp_path):
    f = tmp_path / "corrupt.csv"
    f.write_bytes(bytes(range(256)) * 4)
    units, error = safe_extract(f)
    assert isinstance(units, list)

    from datasentinel_agent.core.enums import ScanStatus
    from datasentinel_agent.core.pipeline import ScanOptions, run_scan
    from datasentinel_agent.storage.database import init_db, make_engine, make_session_factory

    engine = make_engine(tmp_path / "agent_test.db")
    init_db(engine)
    session_factory = make_session_factory(engine)
    options = ScanOptions(profile="standard", paths=[tmp_path], use_presidio=False)
    summary = run_scan(options, session_factory)
    assert summary.status == ScanStatus.COMPLETED


def test_corrupted_xlsx_never_raises(tmp_path):
    pytest.importorskip("openpyxl")
    f = tmp_path / "corrupt.xlsx"
    f.write_bytes(b"not a real zip/xlsx")
    units, error = safe_extract(f)
    assert units == []
    assert error is not None

    from datasentinel_agent.core.enums import ScanStatus
    from datasentinel_agent.core.pipeline import ScanOptions, run_scan
    from datasentinel_agent.storage.database import init_db, make_engine, make_session_factory

    engine = make_engine(tmp_path / "agent_test.db")
    init_db(engine)
    session_factory = make_session_factory(engine)
    options = ScanOptions(profile="standard", paths=[tmp_path], use_presidio=False)
    summary = run_scan(options, session_factory)
    assert summary.status == ScanStatus.COMPLETED


def test_corrupted_pptx_never_raises(tmp_path):
    pytest.importorskip("pptx")
    f = tmp_path / "corrupt.pptx"
    f.write_bytes(b"not a real zip/pptx")
    units, error = safe_extract(f)
    assert units == []
    assert error is not None

    from datasentinel_agent.core.enums import ScanStatus
    from datasentinel_agent.core.pipeline import ScanOptions, run_scan
    from datasentinel_agent.storage.database import init_db, make_engine, make_session_factory

    engine = make_engine(tmp_path / "agent_test.db")
    init_db(engine)
    session_factory = make_session_factory(engine)
    options = ScanOptions(profile="standard", paths=[tmp_path], use_presidio=False)
    summary = run_scan(options, session_factory)
    assert summary.status == ScanStatus.COMPLETED
